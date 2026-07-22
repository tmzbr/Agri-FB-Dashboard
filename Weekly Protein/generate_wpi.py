"""
generate_wpi.py
---------------
Weekly Protein Intake — PPT Generator
Agribusiness, F&B DataHouse · Food & Beverage Team

Usage:
    python generate_wpi.py --week 129 --date "June 26, 2026"

If --date is omitted, today's date is used.
The script downloads the latest .db files from GitHub, generates all charts,
and saves the updated PPT as WPI_#NNN_YYYYMMDD.pptx
"""

import argparse, io, os, shutil, sqlite3, tempfile, time, urllib.request, warnings
from datetime import date
from copy import deepcopy

import numpy as np
import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib as mpl
import matplotlib.ticker
mpl.rcParams.update({
    'font.family':     'sans-serif',
    'font.sans-serif': ['Arial', 'Helvetica', 'DejaVu Sans'],
    'font.size':       6,
    'text.color':      '#000000',
    'axes.labelcolor': '#000000',
    'xtick.color':     '#000000',
    'ytick.color':     '#000000',
})
FONT_SZ   = 6    # standard label / tick size
FONT_SZ_S = 5.5  # small (x-tick labels)
FONT_SZ_L = 6.5  # large (bar value labels)
import matplotlib.patches as mpatches
from matplotlib.lines import Line2D
from scipy.ndimage import gaussian_filter1d
from lxml import etree

from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.oxml.ns import qn
from pptx.enum.shapes import PP_PLACEHOLDER

warnings.filterwarnings('ignore')

# ─────────────────────────────────────────────────────────────────────────────
# CONFIG
# ─────────────────────────────────────────────────────────────────────────────
TEMPLATE   = os.path.join(os.path.dirname(__file__), 'FOOD_AND_BEVERAGE_WPI.pptx')
OUTPUT_DIR = os.path.dirname(__file__)
GITHUB_BASE = "https://raw.githubusercontent.com/tmzbr/Agri-FB-Dashboard/main"

DB_URLS = {
    'beef_bz':    f"{GITHUB_BASE}/Beef/BZ%20Tracker/beef_bz.db",
    'chicken_bz': f"{GITHUB_BASE}/Chicken/BZ%20Tracker/chicken_bz.db",
    'pork_bz':    f"{GITHUB_BASE}/Pork/BZ%20Tracker/pork_bz.db",
    'beef_us':    f"{GITHUB_BASE}/Beef/U.S.%20Tracker/beef.db",
    'chicken_us': f"{GITHUB_BASE}/Chicken/U.S.%20Tracker/chicken.db",
    'pork_us':    f"{GITHUB_BASE}/Pork/U.S.%20Tracker/pork_us.db",
}

# Chart OLE positions per slide (inches) — from template inspection
# Each tuple: (left, top, width, height)
CHART_POS = {
    3: [  # Beef BZ
        (0.47, 1.48, 5.74, 2.00),   # top-left:    Spread series
        (7.22, 1.48, 5.70, 1.98),   # top-right:   Vol & Price bars
        (0.57, 4.61, 5.79, 1.95),   # bottom-left: Historical seasonal
        (7.15, 4.60, 5.78, 1.96),   # bottom-right: Price variation
    ],
    4: [  # Chicken BZ
        (0.40, 1.52, 5.76, 1.99),
        (7.16, 1.52, 5.77, 1.96),
        (0.49, 4.64, 5.95, 1.96),
        (7.16, 4.63, 5.86, 1.98),
    ],
    5: [  # Pork BZ
        (0.42, 1.47, 5.70, 2.04),
        (7.23, 1.47, 5.74, 1.95),
        (0.46, 4.57, 5.78, 1.96),
        (7.20, 4.57, 5.74, 1.96),
    ],
    6: [  # US Domestic
        (0.58, 1.45, 3.87, 2.01),   # Beef spread (recent)
        (4.84, 1.45, 3.94, 2.01),   # Poultry spread (recent)
        (9.35, 1.45, 3.81, 2.01),   # Pork spread (recent)
        (0.61, 4.66, 3.84, 2.01),   # Beef historical
        (4.92, 4.66, 3.85, 2.00),   # Poultry historical
        (9.35, 4.66, 3.78, 2.01),   # Pork historical
    ],
}

# ─────────────────────────────────────────────────────────────────────────────
# PALETTE
# ─────────────────────────────────────────────────────────────────────────────
BG     = '#E0E0E0'
ORANGE = '#FF5500'
BLACK  = '#000000'
GRAY   = '#888888'
BLUE   = '#0070C0'
LGRAY  = '#CCCCCC'

# ─────────────────────────────────────────────────────────────────────────────
# STEP 1 — Download DBs
# ─────────────────────────────────────────────────────────────────────────────
def download_dbs(tmpdir):
    """
    Download DBs from GitHub. If a local override file exists next to this
    script (e.g. pork_us.db), use it instead — useful when the GitHub DB
    has not yet been updated with corrected data.
    """
    script_dir = os.path.dirname(os.path.abspath(__file__))
    dbs = {}
    for name, url in DB_URLS.items():
        dest = os.path.join(tmpdir, f'{name}.db')
        # Check for local override
        local = os.path.join(script_dir, f'{name}.db')
        if os.path.exists(local):
            import shutil
            shutil.copy(local, dest)
            dbs[name] = dest
            print(f"  {name}.db  [LOCAL OVERRIDE] ({os.path.getsize(dest)//1024}KB)")
            continue
        print(f"  Downloading {name}.db ...", end=' ')
        try:
            urllib.request.urlretrieve(url, dest)
            dbs[name] = dest
            print(f"OK ({os.path.getsize(dest)//1024}KB)")
        except Exception as ex:
            print(f"FAILED: {ex}")
            dbs[name] = None
        time.sleep(0.3)
    return dbs

# ─────────────────────────────────────────────────────────────────────────────
# STEP 2 — Load data
# ─────────────────────────────────────────────────────────────────────────────
def load_bz(path):
    """Load BZ tracker: returns (df_weekly, df_monthly)."""
    con = sqlite3.connect(path)
    dfw = pd.read_sql("SELECT * FROM weekly  ORDER BY start_date", con)
    dfm = pd.read_sql("SELECT * FROM monthly ORDER BY period",     con)
    con.close()
    dfw['start_date'] = pd.to_datetime(dfw['start_date'])
    dfm['period']     = pd.to_datetime(dfm['period'])
    return dfw, dfm

# ─────────────────────────────────────────────────────────────────────────────
# STEP 3 — Chart helpers
# ─────────────────────────────────────────────────────────────────────────────
import matplotlib.ticker

FS = 6  # uniform font size for all chart labels, ticks, legends

matplotlib.rcParams.update({
    'font.family':     'sans-serif',
    'font.sans-serif': ['Arial', 'Helvetica', 'DejaVu Sans'],
    'font.size':       FS,
    'text.color':      '#000000',
    'axes.labelcolor': '#000000',
    'xtick.color':     '#000000',
    'ytick.color':     '#000000',
})

def nice_scale(lo, hi, n=5):
    rng = hi - lo
    if rng == 0: rng = 1
    raw = rng / n
    mag = 10 ** np.floor(np.log10(raw + 1e-9))
    for step in [1, 2, 2.5, 5, 10]:
        if raw <= step * mag:
            s = step * mag; break
    return np.floor(lo / s) * s, np.ceil(hi / s) * s, s

def save_fig(fig, path):
    # Preserve figure facecolor (transparent for line charts, BG for bar charts)
    fc = fig.get_facecolor()
    fig.savefig(path, format='png', dpi=200, bbox_inches='tight', facecolor=fc)
    plt.close(fig)

def chart_title(ax, title, fs=8):
    pass  # titles shown in PPTX template

def smooth(y, sigma=2.0):
    arr = np.asarray(y, dtype=float)
    return gaussian_filter1d(arr, sigma=sigma) if len(arr) >= 4 else arr

def setup(figsize, bottom_only=False):
    """Transparent figure, clean axes, bottom (+optional left) spine only."""
    fig, ax = plt.subplots(figsize=figsize)
    fig.patch.set_alpha(0.0)
    ax.set_facecolor('none')
    ax.grid(False)
    ax.spines[:].set_visible(False)
    ax.tick_params(which='both', length=0)
    ax.spines['bottom'].set_visible(True)
    ax.spines['bottom'].set_color(BLACK)
    ax.spines['bottom'].set_linewidth(0.7)
    if not bottom_only:
        ax.spines['left'].set_visible(True)
        ax.spines['left'].set_color(BLACK)
        ax.spines['left'].set_linewidth(0.7)
    return fig, ax

def week_avg(dfw, col, yr, mo):
    r = dfw[(dfw['start_date'].dt.year==yr) & (dfw['start_date'].dt.month==mo)]
    return float(r[col].mean()) if len(r) else 0.0


# ─────────────────────────────────────────────────────────────────────────────
# BZ CHARTS (slides 3–5)
# ─────────────────────────────────────────────────────────────────────────────
def bz_chart1_spread(dfw, dfm, protein_label, out):
    """Spread time series — smooth, eixo bottom+left."""
    df = dfw[dfw['start_date'] >= '2022-09-01'].copy()
    avg_lo = dfm[(dfm['year']>=2015)&(dfm['year']<=2019)]['spread'].mean()
    avg_hi = dfm[(dfm['year']>=2020)&(dfm['year']<=2025)]['spread'].mean()

    all_vals = list(df['spread']) + [avg_lo, avg_hi]
    lo, hi, s = nice_scale(min(all_vals)*0.97, max(all_vals)*1.03)

    fig, ax = setup((5.5, 2.0))
    y_sm = smooth(df['spread'].values, sigma=1.5)
    ax.plot(df['start_date'], y_sm, color=BLACK, lw=0.9, zorder=3)
    ax.axhline(avg_lo, color=ORANGE, lw=0.9, ls='--', zorder=2)
    ax.axhline(avg_hi, color=GRAY,   lw=0.9, ls='--', zorder=2)
    # Last point orange
    ax.plot(df['start_date'].iloc[-1], y_sm[-1], 'o', color=ORANGE, ms=3.5, zorder=5)

    ax.set_ylim(lo, hi)
    ax.set_xlim(df['start_date'].iloc[0], df['start_date'].iloc[-1])
    ax.set_yticks(np.arange(lo, hi+s*0.01, s))
    ax.yaxis.set_tick_params(labelsize=FS)
    xt = df['start_date'].iloc[::6]
    ax.set_xticks(xt)
    ax.set_xticklabels([d.strftime('%b-%y') for d in xt], rotation=90, fontsize=FS)

    handles = [
        Line2D([0],[0], color=BLACK,  lw=0.9, label=f'{protein_label} Exports Spread'),
        Line2D([0],[0], color=ORANGE, lw=0.9, ls='--', label='2015-2019 Average'),
        Line2D([0],[0], color=GRAY,   lw=0.9, ls='--', label='2020-2025 Average'),
    ]
    ax.legend(handles=handles, fontsize=FS, frameon=False,
              loc='lower center', ncol=3, bbox_to_anchor=(0.5, -1.20))
    fig.subplots_adjust(left=0.10, right=0.98, top=0.96, bottom=0.74)
    save_fig(fig, out)
def bz_chart2_vol_price(dfw, protein_label, out):
    """Vol bars + price dots. This Week and current month orange; rest black. Transparent bg."""
    latest  = dfw['start_date'].max()
    cur_yr, cur_mo = latest.year, latest.month
    prev_mo = cur_mo-1 if cur_mo>1 else 12
    prev_yr = cur_yr  if cur_mo>1 else cur_yr-1

    lw_row = dfw.iloc[-2]; tw_row = dfw.iloc[-1]
    labels = ['Last Week','This Week',
              f'{latest.strftime("%b")}-{cur_yr-1}',
              f'{pd.Timestamp(prev_yr,prev_mo,1).strftime("%b")}-{prev_yr}',
              f'{latest.strftime("%b")}-{cur_yr}']
    # This Week (1) and current month (4) = orange; rest = black
    bar_colors = [BLACK, ORANGE, BLACK, BLACK, ORANGE]
    vols   = [lw_row['vol_tons_daily']/1000, tw_row['vol_tons_daily']/1000,
              week_avg(dfw,'vol_tons_daily',cur_yr-1,cur_mo)/1000,
              week_avg(dfw,'vol_tons_daily',prev_yr,prev_mo)/1000,
              week_avg(dfw,'vol_tons_daily',cur_yr,cur_mo)/1000]
    prices = [lw_row['secex_usd_kg'], tw_row['secex_usd_kg'],
              week_avg(dfw,'secex_usd_kg',cur_yr-1,cur_mo),
              week_avg(dfw,'secex_usd_kg',prev_yr,prev_mo),
              week_avg(dfw,'secex_usd_kg',cur_yr,cur_mo)]

    v_max    = max(vols) if vols else 1
    ylim_top = v_max * 1.85

    fig, ax = plt.subplots(figsize=(4.5, 2.0))
    fig.patch.set_alpha(0.0)
    ax.set_facecolor('none')
    ax.grid(False)
    ax.spines[:].set_visible(False)
    ax.spines['bottom'].set_visible(True)
    ax.spines['bottom'].set_color(BLACK)
    ax.spines['bottom'].set_linewidth(0.7)
    ax.spines['left'].set_visible(True)
    ax.spines['left'].set_color(BLACK)
    ax.spines['left'].set_linewidth(0.7)
    ax.tick_params(which='both', length=0)

    x = np.arange(len(labels))
    for xi, (v, col) in enumerate(zip(vols, bar_colors)):
        ax.bar(xi, v, color=col, width=0.5, zorder=3)

    # Dots and labels on same axis — dots fixed above all bars
    dot_y = v_max * 1.30
    lbl_y = v_max * 1.43
    valid_p = [p for p in prices if p]
    p_max   = max(valid_p) if valid_p else 1

    for xi, (v, p) in enumerate(zip(vols, prices)):
        ax.plot(xi, dot_y, 'o', color=GRAY, ms=5, zorder=6, clip_on=False)
        if p:
            ax.text(xi, lbl_y, f'{p:.1f}',
                    ha='center', va='bottom', fontsize=FS, color=BLACK)
        if v:
            ax.text(xi, v + v_max*0.04, f'{v:.1f}',
                    ha='center', va='bottom', fontsize=FS, color=BLACK, fontweight='bold')

    ax.set_xticks(x); ax.set_xticklabels(labels, fontsize=FS)
    ax.set_ylim(0, ylim_top); ax.yaxis.set_visible(False)

    handles = [mpatches.Patch(color=BLACK, label="Daily Volume ('000 tons)"),
               Line2D([0],[0], marker='o', color='none', markerfacecolor=GRAY,
                      markersize=4, label='Price (USD/Kg)')]
    ax.legend(handles=handles, fontsize=FS, frameon=False,
              loc='lower center', ncol=2, bbox_to_anchor=(0.5, -0.55))
    fig.subplots_adjust(left=0.02, right=0.98, top=0.96, bottom=0.42)
    save_fig(fig, out)
def bz_chart3_historical(dfw, dfm, protein_label, hist_start_year, out):
    """Seasonal spreads — smooth year lines, 2026 only up to last available month."""
    dfm['month_n'] = dfm['month']
    grp = dfm[(dfm['year']>=hist_start_year)&(dfm['year']<=2023)].groupby('month_n')['spread']
    hmin = grp.min(); hmax = grp.max()
    months = np.arange(1,13)

    fig, ax = setup((5.5, 2.0))
    ax.fill_between(months, hmin.reindex(months), hmax.reindex(months),
                    color=GRAY, alpha=0.25, zorder=1)

    for yr, col, lw in [(2024, BLUE, 1.0), (2025, BLACK, 1.1), (2026, ORANGE, 1.3)]:
        s = dfm[dfm['year']==yr].set_index('month_n')['spread']
        if len(s) == 0: continue
        last_mo = int(s.index.max())
        # Fill internal gaps but never extrapolate beyond last real month
        s_range = s.reindex(range(1, last_mo+1)).interpolate()
        y_sm = smooth(s_range.values, sigma=1.2)
        ax.plot(range(1, last_mo+1), y_sm, color=col, lw=lw, zorder=3, label=str(yr))

    all_vals = dfm[dfm['year']>=hist_start_year]['spread']
    lo, hi, s = nice_scale(all_vals.min()*0.95, all_vals.max()*1.05)
    ax.set_ylim(lo, hi); ax.set_xlim(1, 12); ax.set_xticks(months)
    ax.set_xticklabels(['Jan','Feb','Mar','Apr','May','Jun',
                        'Jul','Aug','Sep','Oct','Nov','Dec'], fontsize=FS)
    ax.set_yticks(np.arange(lo, hi+s*0.01, s))
    ax.yaxis.set_tick_params(labelsize=FS)

    handles = [mpatches.Patch(color=GRAY, alpha=0.35, label=f'{hist_start_year}-2023'),
               Line2D([0],[0], color=BLUE,   lw=1.0, label='2024'),
               Line2D([0],[0], color=BLACK,  lw=1.1, label='2025'),
               Line2D([0],[0], color=ORANGE, lw=1.3, label='2026')]
    ax.legend(handles=handles, fontsize=FS, frameon=False,
              loc='lower center', ncol=4, bbox_to_anchor=(0.5, -0.40))
    fig.subplots_adjust(left=0.10, right=0.98, top=0.96, bottom=0.34)
    save_fig(fig, out)
def bz_chart4_price_var(dfw, out):
    """Paired USD/BRL bars. LW/Jun-YA/PrevM=hatched, TW/CurM=solid. BG bg for hatch."""
    latest  = dfw['start_date'].max()
    cur_yr, cur_mo = latest.year, latest.month
    prev_mo = cur_mo-1 if cur_mo>1 else 12
    prev_yr = cur_yr  if cur_mo>1 else cur_yr-1

    lw_row = dfw.iloc[-2]; tw_row = dfw.iloc[-1]
    labels = ['Last Week','This Week',
              f'{latest.strftime("%b")}-{cur_yr-1}',
              f'{pd.Timestamp(prev_yr,prev_mo,1).strftime("%b")}-{prev_yr}',
              f'{latest.strftime("%b")}-{cur_yr}']
    # solid[i]=True → no hatch (This Week=1, current month=4)
    solid = [False, True, False, False, True]
    usd = [lw_row['secex_usd_kg'], tw_row['secex_usd_kg'],
           week_avg(dfw,'secex_usd_kg',cur_yr-1,cur_mo),
           week_avg(dfw,'secex_usd_kg',prev_yr,prev_mo),
           week_avg(dfw,'secex_usd_kg',cur_yr,cur_mo)]
    brl = [lw_row['secex_brl_kg'], tw_row['secex_brl_kg'],
           week_avg(dfw,'secex_brl_kg',cur_yr-1,cur_mo),
           week_avg(dfw,'secex_brl_kg',prev_yr,prev_mo),
           week_avg(dfw,'secex_brl_kg',cur_yr,cur_mo)]
    mx = max(brl) if brl else 1

    fig, ax = plt.subplots(figsize=(4.5, 2.0))
    fig.patch.set_alpha(0.0)
    ax.set_facecolor('none')
    ax.grid(False)
    ax.spines[:].set_visible(False)
    ax.spines['bottom'].set_visible(True)
    ax.spines['bottom'].set_color(BLACK)
    ax.spines['bottom'].set_linewidth(0.7)
    ax.tick_params(which='both', length=0)

    x = np.arange(len(labels)); w = 0.35
    for xi, sol in enumerate(solid):
        if sol:
            # Solid: no hatch, invisible edge
            ax.bar(xi-w/2, usd[xi], w, color='#1A1A1A',
                   edgecolor=BG, linewidth=0.3, zorder=3,
                   label='Price (USD/Kg)' if xi==0 else None)
            ax.bar(xi+w/2, brl[xi], w, color=ORANGE,
                   edgecolor=BG, linewidth=0.3, zorder=3,
                   label='Price (BRL/Kg)' if xi==0 else None)
        else:
            # Hatched: white edgecolor so hatch lines are visible
            ax.bar(xi-w/2, usd[xi], w, color='#1A1A1A', hatch='//',
                   edgecolor='white', linewidth=0.8, zorder=3,
                   label='Price (USD/Kg)' if xi==0 else None)
            ax.bar(xi+w/2, brl[xi], w, color=ORANGE, hatch='//',
                   edgecolor='white', linewidth=0.8, zorder=3,
                   label='Price (BRL/Kg)' if xi==0 else None)

    for xi, (u, b) in enumerate(zip(usd, brl)):
        if u: ax.text(xi-w/2, u+mx*0.02, f'{u:.1f}',
                      ha='center', va='bottom', fontsize=FS, color=BLACK, fontweight='bold')
        if b: ax.text(xi+w/2, b+mx*0.02, f'{b:.1f}',
                      ha='center', va='bottom', fontsize=FS, color=BLACK, fontweight='bold')

    ax.set_xticks(x); ax.set_xticklabels(labels, fontsize=FS)
    ax.set_ylim(0, mx*1.30); ax.yaxis.set_visible(False)
    ax.legend(fontsize=FS, frameon=False, loc='lower center',
              ncol=2, bbox_to_anchor=(0.5, -0.28))
    fig.subplots_adjust(left=0.02, right=0.98, top=0.96, bottom=0.26)
    save_fig(fig, out)
def us_spread_recent(df_w, spread_col, avg_cols, label, date_col, out,
                     avg_labels=('Average 2015-2020','Average 2021-2025')):
    """Recent US spread — smooth, transparent, bottom+left axes, last point orange."""
    df_w = df_w.copy()
    df_w[date_col] = pd.to_datetime(df_w[date_col])
    df_w = df_w.dropna(subset=[spread_col]).sort_values(date_col)
    cutoff = df_w[date_col].max() - pd.DateOffset(months=13)
    recent = df_w[df_w[date_col] >= cutoff]
    if recent.empty: recent = df_w.tail(30)

    avg_lo = df_w[df_w[date_col].dt.year.isin(range(avg_cols[0][0], avg_cols[0][1]+1))][spread_col].mean()
    avg_hi = df_w[df_w[date_col].dt.year.isin(range(avg_cols[1][0], avg_cols[1][1]+1))][spread_col].mean()

    all_vals = list(recent[spread_col].dropna()) + [avg_lo, avg_hi]
    lo, hi, s = nice_scale(min(all_vals)*0.97, max(all_vals)*1.03)

    fig, ax = setup((3.7, 2.0))
    y_sm = smooth(recent[spread_col].values, sigma=1.5)
    ax.plot(recent[date_col], y_sm, color=BLACK, lw=0.9, zorder=3)
    ax.axhline(avg_lo, color=ORANGE, lw=0.9, ls='--', zorder=2)
    ax.axhline(avg_hi, color=GRAY,   lw=0.9, ls='--', zorder=2)
    ax.plot(recent[date_col].iloc[-1], y_sm[-1], 'o', color=ORANGE, ms=3.5, zorder=5)

    ax.set_ylim(lo, hi); ax.set_xlim(recent[date_col].iloc[0], recent[date_col].iloc[-1])
    ax.set_yticks(np.arange(lo, hi+s*0.01, s))
    ax.yaxis.set_tick_params(labelsize=FS)
    xt = recent[date_col].iloc[::5]
    ax.set_xticks(xt)
    ax.set_xticklabels([d.strftime('%b-%y') for d in xt], rotation=90, fontsize=FS)

    handles = [Line2D([0],[0], color=BLACK,  lw=0.9, label=label),
               Line2D([0],[0], color=ORANGE, lw=0.9, ls='--', label=avg_labels[0]),
               Line2D([0],[0], color=GRAY,   lw=0.9, ls='--', label=avg_labels[1])]
    ax.legend(handles=handles, fontsize=FS, frameon=False,
              loc='lower center', ncol=3, bbox_to_anchor=(0.5, -1.20))
    fig.subplots_adjust(left=0.12, right=0.98, top=0.96, bottom=0.74)
    save_fig(fig, out)
def us_spread_historical(df_w, spread_col, date_col, hist_start, label, out,
                         exclude_years=None):
    """Historical US spread — smooth, transparent, bottom+left axes."""
    df_w = df_w.copy()
    df_w[date_col] = pd.to_datetime(df_w[date_col])
    df_w = df_w.dropna(subset=[spread_col]).sort_values(date_col)
    df_w['month'] = df_w[date_col].dt.month
    df_w['year']  = df_w[date_col].dt.year

    mask = (df_w['year'] >= hist_start) & (df_w['year'] <= 2023)
    if exclude_years:
        mask = mask & (~df_w['year'].isin(exclude_years))
    grp  = df_w[mask].groupby('month')[spread_col]
    hmin = grp.min(); hmax = grp.max()
    months = np.arange(1,13)

    fig, ax = setup((3.7, 2.0))
    ax.fill_between(months, hmin.reindex(months), hmax.reindex(months),
                    color=GRAY, alpha=0.25, zorder=1)

    for yr, col, lw in [(2024, BLUE, 1.0), (2025, BLACK, 1.1), (2026, ORANGE, 1.3)]:
        s_raw = df_w[df_w['year']==yr].groupby('month')[spread_col].mean()
        if len(s_raw) < 2: continue
        last_mo = int(s_raw.index.max())
        s_full  = s_raw.reindex(range(1, last_mo+1)).interpolate()
        y_sm    = smooth(s_full.values, sigma=1.2)
        ax.plot(range(1, last_mo+1), y_sm, color=col, lw=lw, zorder=3, label=str(yr))

    scale_data = df_w[mask | df_w['year'].isin([2024,2025,2026])][spread_col]
    lo, hi, s  = nice_scale(scale_data.min()*0.95, scale_data.max()*1.05)
    ax.set_ylim(lo, hi); ax.set_xlim(1,12); ax.set_xticks(months)
    ax.set_xticklabels(['Jan','Feb','Mar','Apr','May','Jun',
                        'Jul','Aug','Sep','Oct','Nov','Dec'], fontsize=FS)
    ax.set_yticks(np.arange(lo, hi+s*0.01, s))
    ax.yaxis.set_major_formatter(matplotlib.ticker.FormatStrFormatter('%.1f'))
    ax.yaxis.set_tick_params(labelsize=FS)

    handles = [mpatches.Patch(color=GRAY, alpha=0.35, label=f'{hist_start}-2023'),
               Line2D([0],[0], color=BLUE,   lw=1.0, label='2024'),
               Line2D([0],[0], color=BLACK,  lw=1.1, label='2025'),
               Line2D([0],[0], color=ORANGE, lw=1.3, label='2026')]
    ax.legend(handles=handles, fontsize=FS, frameon=False,
              loc='lower center', ncol=4, bbox_to_anchor=(0.5, -0.30))
    fig.subplots_adjust(left=0.12, right=0.98, top=0.96, bottom=0.28)
    save_fig(fig, out)
def replace_ole_with_image(slide, ole_idx, img_path, pos):
    """Remove OLE object at index and insert PNG at same position."""
    left, top, w, h = [Inches(v) for v in pos]
    # Remove the OLE object
    shape = slide.shapes[ole_idx]
    sp_el = shape._element
    sp_el.getparent().remove(sp_el)
    # Insert picture
    slide.shapes.add_picture(img_path, left, top, w, h)

def update_slide_charts(slide, chart_paths, positions):
    """Replace all OLE objects on a slide with chart images."""
    # Collect OLE objects by position order (left → right, top → bottom)
    ole_shapes = [(i, sh) for i, sh in enumerate(slide.shapes)
                  if sh.shape_type == 10]  # LINKED_OLE_OBJECT
    # Sort by (top, left) to match our position list
    ole_shapes.sort(key=lambda x: (round(x[1].top/914400, 1), round(x[1].left/914400, 1)))

    # Remove all OLE elements from XML first
    elements_to_remove = [sh._element for _, sh in ole_shapes]
    for el in elements_to_remove:
        el.getparent().remove(el)

    # Now add picture at each position
    for img_path, pos in zip(chart_paths, positions):
        left, top, w, h = [Inches(v) for v in pos]
        slide.shapes.add_picture(img_path, left, top, w, h)

# ─────────────────────────────────────────────────────────────────────────────
# STEP 5 — Update text (slide 1: week number and date)
# ─────────────────────────────────────────────────────────────────────────────
def update_cover(slide, week_number, report_date_str):
    for shape in slide.shapes:
        if not hasattr(shape, 'text'): continue
        if '#xx' in shape.text:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if '#xx' in run.text:
                        run.text = run.text.replace('#xx', f'#{week_number}')
        if 'June xx' in shape.text or 'xx, 2026' in shape.text.lower() or shape.name == 'Text Box 10':
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    # Replace date placeholder
                    import re
                    run.text = re.sub(r'[A-Z][a-z]+ xx, \d{4}', report_date_str, run.text)

# ─────────────────────────────────────────────────────────────────────────────
# MAIN
# ─────────────────────────────────────────────────────────────────────────────
def main():
    parser = argparse.ArgumentParser(description='Generate Weekly Protein Intake PPT')
    parser.add_argument('--week', type=int, default=None, help='Week number (e.g. 130). Auto-detected from metadata if omitted.')
    parser.add_argument('--date', type=str, default=None,
                        help='Report date string, e.g. "July 3, 2026" (default: today)')
    parser.add_argument('--template', type=str, default=TEMPLATE,
                        help='Path to PPTX template')
    args = parser.parse_args()

    report_date = args.date or date.today().strftime('%B %-d, %Y')

    # ── Auto-detect week number ──────────────────────────────────────────────
    week_number = args.week
    if week_number is None:
        # Try to read from wpi_meta.json next to the script
        meta_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'wpi_meta.json')
        if os.path.exists(meta_path):
            import json
            with open(meta_path) as f:
                meta = json.load(f)
            week_number = meta.get('last_week', 128) + 1
            print(f"  [AUTO] Week number: #{week_number} (last was #{week_number-1} per wpi_meta.json)")
        else:
            week_number = 130
            print(f"  [AUTO] wpi_meta.json not found — defaulting to #{week_number}. Use --week to specify.")
    args.week = week_number
    print(f"\n{'='*55}")
    print(f"  Weekly Protein Intake #{args.week}  |  {report_date}")
    print(f"{'='*55}\n")

    tmpdir = tempfile.mkdtemp()
    chart_dir = os.path.join(tmpdir, 'charts')
    os.makedirs(chart_dir)

    try:
        # ── Download DBs ──────────────────────────────────────────────────
        print("1. Downloading databases from GitHub...")
        dbs = download_dbs(tmpdir)

        # ── Generate BZ charts ───────────────────────────────────────────
        print("\n2. Generating Brazil export charts...")
        bz_configs = [
            ('beef_bz',    'Beef',    3, 2017),
            ('chicken_bz', 'Poultry', 4, 2017),
            ('pork_bz',    'Pork',    5, 2017),
        ]
        bz_chart_files = {}
        for db_key, label, slide_n, hist_yr in bz_configs:
            if not dbs.get(db_key):
                print(f"  SKIP {label} (no db)")
                continue
            dfw, dfm = load_bz(dbs[db_key])
            charts = []
            for fn, args_extra, suffix in [
                (bz_chart1_spread,    (label,),          'spread'),
                (bz_chart2_vol_price, (label,),          'vol'),
                (bz_chart3_historical,(label, hist_yr),  'hist'),
                (bz_chart4_price_var, (),                'price'),
            ]:
                out = os.path.join(chart_dir, f's{slide_n}_{suffix}.png')
                fn(dfw, dfm, *args_extra, out) if 'dfm' in fn.__code__.co_varnames else fn(dfw, *args_extra, out)
                charts.append(out)
                print(f"  ✅ Slide {slide_n} — {suffix}")
            bz_chart_files[slide_n] = charts

        # Pork historical starts 2017 (data available from 2017)

        # ── Generate U.S. Domestic Spread charts ──────────────────────────
        # Formulas verified against Excel U.S. Spreads - Weekly sheet:
        #   Beef:    choice / ct150_steer
        #   Pork:    carcass / hog_price
        #   Poultry: (weighted_carcass_USD/kg) / (fc_spot/100 lagged 2Q * 3)
        #            where carcass = (breast*29%+wings*10%+leg_qtrs*41%+tenders*5%)/85% /0.4536/100
        print("\n3. Generating U.S. Domestic Spread charts...")
        us_charts = []

        # ── BEEF US ──────────────────────────────────────────────────────────
        if dbs.get('beef_us'):
            con = sqlite3.connect(dbs['beef_us'])
            df_bw = pd.read_sql("SELECT * FROM beef_weekly ORDER BY week_ending", con)
            con.close()
            df_bw['spread'] = df_bw['choice'] / df_bw['ct150_steer']
            out = os.path.join(chart_dir, 's6_beef_recent.png')
            us_spread_recent(df_bw, 'spread', ((2015,2020),(2021,2025)),
                             'Beef Spread', 'week_ending', out)
            us_charts.append(out); print("  ✅ Slide 6 — Beef Spread (recent)")
            out = os.path.join(chart_dir, 's6_beef_hist.png')
            us_spread_historical(df_bw, 'spread', 'week_ending', 2017,
                                 'U.S. Beef Spreads', out, exclude_years=[2020])
            us_charts.append(out); print("  ✅ Slide 6 — Beef Spread (historical)")
        else:
            us_charts += [None, None]

        # ── POULTRY US ────────────────────────────────────────────────────────
        # Source: chicken_us.db
        # Spread = PoultCarcass(USD/kg) / (FeedGrain_2m(USD/kg) * 3)  [FCR=3]
        # Carcass blend: (breast*29% + wings*10% + leg_qtrs*41% + tenders*5%) / 85% /0.4536/100
        # Feed: fc_spot/100 USD/kg (quarterly), lagged 2 quarters
        if dbs.get('chicken_us'):
            # Fallback: DB with approximate formula
            con = sqlite3.connect(dbs['chicken_us'])
            dfw_c = pd.read_sql("SELECT * FROM weekly  ORDER BY report_date", con)
            dfq_c = pd.read_sql("SELECT * FROM quarterly ORDER BY year_q",  con)
            con.close()
            for col in ['breast','wings','leg_qtrs','tenders']:
                dfw_c[col] = pd.to_numeric(dfw_c[col], errors='coerce')
            dfw_c['report_date'] = pd.to_datetime(dfw_c['report_date'])
            dfw_c['carcass_usd_kg'] = (
                dfw_c['breast']*(0.29/0.85)+dfw_c['wings']*(0.10/0.85)+
                dfw_c['leg_qtrs']*(0.41/0.85)+dfw_c['tenders']*(0.05/0.85)
            ) / 0.4536 / 100
            dfq_c['year_n'] = dfq_c['year_q'].astype(str).str[:4].astype(int)
            dfq_c['q_n']    = dfq_c['year_q'].astype(str).str[4:].astype(int)
            dfq_c['date']   = pd.to_datetime(dfq_c['year_n'].astype(str)+'-'+((dfq_c['q_n']-1)*3+1).astype(str).str.zfill(2)+'-01')
            dfq_c['feed_lag'] = (dfq_c['fc_spot']/100).shift(2)
            dfq_sel = dfq_c[['date','feed_lag']].set_index('date')
            feed_w = dfq_sel.reindex(dfq_sel.index.union(dfw_c['report_date'])).ffill().reindex(dfw_c['report_date'])
            dfw_c['feed_lag'] = feed_w['feed_lag'].values
            dfw_c['spread']   = dfw_c['carcass_usd_kg']/(dfw_c['feed_lag']*3)
            df_chk = dfw_c.dropna(subset=['spread']).copy()
            out = os.path.join(chart_dir, 's6_chk_recent.png')
            us_spread_recent(df_chk, 'spread', ((2015,2020),(2021,2025)),
                             'Poultry Spread (2-Month Lag)', 'report_date', out)
            us_charts.insert(1, out); print("  ✅ Slide 6 — Poultry Spread (recent) [DB fallback]")
            out = os.path.join(chart_dir, 's6_chk_hist.png')
            us_spread_historical(df_chk, 'spread', 'report_date', 2017,
                                 'U.S. Poultry Spreads (2-Month Lag)', out)
            us_charts.insert(3, out); print("  ✅ Slide 6 — Poultry Spread (historical) [DB fallback]")
        else:
            us_charts.insert(1, None); us_charts.insert(3, None)

        # ── PORK US ──────────────────────────────────────────────────────────
        # Source: pork_us.db (rebuilt with BOXPC185/HOGSNATL series = LM_PK601/LM_HG203)
        # Spread = carcass / hog_price  (both USD/cwt, same basis as Bloomberg)
        if dbs.get('pork_us'):
            con = sqlite3.connect(dbs['pork_us'])
            df_pw = pd.read_sql("SELECT * FROM weekly ORDER BY report_date", con)
            con.close()
            df_pw['carcass']   = pd.to_numeric(df_pw['carcass'],   errors='coerce')
            df_pw['hog_price'] = pd.to_numeric(df_pw['hog_price'], errors='coerce')
            df_pw['spread']    = df_pw['carcass'] / df_pw['hog_price']
            df_pw_clean        = df_pw.dropna(subset=['spread'])
            out = os.path.join(chart_dir, 's6_pork_recent.png')
            us_spread_recent(df_pw_clean, 'spread', ((2015,2020),(2021,2025)),
                             'Pork Spread (Non-Integrated)', 'report_date', out)
            us_charts.insert(2, out); print("  ✅ Slide 6 — Pork Spread (recent)")
            out = os.path.join(chart_dir, 's6_pork_hist.png')
            us_spread_historical(df_pw_clean, 'spread', 'report_date', 2017,
                                 'U.S. Pork Spreads (Non-Integrated)', out, exclude_years=[2020])
            us_charts.insert(5, out); print("  ✅ Slide 6 — Pork Spread (historical)")
        else:
            us_charts.insert(2, None); us_charts.insert(5, None)

        while len(us_charts) < 6: us_charts.append(None)

        # ── Edit PPTX ────────────────────────────────────────────────────
        print("\n4. Updating PPTX template...")
        prs = Presentation(args.template)

        # Slide 1 — update week number and date
        update_cover(prs.slides[0], args.week, report_date)
        print(f"  ✅ Slide 1 — cover updated (#{args.week}, {report_date})")

        # Slides 3-5 — BZ charts
        for slide_n in [3, 4, 5]:
            if slide_n not in bz_chart_files: continue
            slide = prs.slides[slide_n - 1]
            update_slide_charts(slide, bz_chart_files[slide_n], CHART_POS[slide_n])
            print(f"  ✅ Slide {slide_n} — charts replaced")

        # Slide 6 — US charts
        valid_us = [c for c in us_charts if c]
        if valid_us:
            slide6 = prs.slides[5]
            update_slide_charts(slide6, us_charts, CHART_POS[6])
            print(f"  ✅ Slide 6 — charts replaced")

        # ── Save ─────────────────────────────────────────────────────────
        today_str = date.today().strftime('%Y%m%d')
        out_name = f"WPI_{args.week:03d}_{today_str}.pptx"
        out_path = os.path.join(OUTPUT_DIR, out_name)
        prs.save(out_path)
        # Save meta for next auto-increment
        meta_path = os.path.join(OUTPUT_DIR, 'wpi_meta.json')
        import json
        with open(meta_path, 'w') as mf:
            json.dump({'last_week': args.week, 'last_date': report_date,
                       'last_file': out_name, 'generated_at': date.today().isoformat()}, mf, indent=2)

        print(f"\n{'='*55}")
        print(f"  ✅ Saved: {out_path}")
        print(f"  ✅ Meta:  {meta_path}")
        print(f"{'='*55}\n")

    finally:
        shutil.rmtree(tmpdir, ignore_errors=True)

if __name__ == '__main__':
    main()
